"""
Document conversion services (PDF <-> Word, Excel, PowerPoint).
"""

import io
import os
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Tuple

from django.conf import settings
import fitz  # PyMuPDF
from PIL import Image


class DocumentConversionService:
    """Service for converting between PDF and other document formats."""

    def __init__(self):
        self.temp_dir = settings.PDF_TEMP_DIR
        self.temp_dir.mkdir(parents=True, exist_ok=True)

    def _get_temp_path(self, suffix: str) -> Path:
        """Generate a temporary file path."""
        import uuid
        return self.temp_dir / f"{uuid.uuid4()}{suffix}"

    def _cleanup_files(self, *paths):
        """Clean up temporary files."""
        for path in paths:
            try:
                if path and Path(path).exists():
                    Path(path).unlink()
            except Exception:
                pass

    # ==================== PDF TO WORD ====================
    def pdf_to_word(self, pdf_file: io.BytesIO) -> Tuple[bytes, str]:
        """
        Convert PDF to Word document.
        Uses text extraction + python-docx for basic conversion.
        
        For production, consider using LibreOffice or a cloud API.
        """
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        pdf_doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        for page_num, page in enumerate(pdf_doc):
            # Extract text
            text = page.get_text()
            
            if text.strip():
                # Add page content
                for paragraph in text.split('\n\n'):
                    if paragraph.strip():
                        p = doc.add_paragraph(paragraph.strip())
            
            # Add page break between pages (except last)
            if page_num < len(pdf_doc) - 1:
                doc.add_page_break()
        
        pdf_doc.close()
        
        # Save to bytes
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        
        filename = f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        return output.read(), filename

    # ==================== WORD TO PDF ====================
    def word_to_pdf(self, word_file: io.BytesIO, original_filename: str = "document.docx") -> Tuple[bytes, str]:
        """
        Convert Word document to PDF.
        Uses LibreOffice for conversion.
        """
        input_path = self._get_temp_path('.docx')
        output_dir = self.temp_dir
        
        try:
            # Save input file
            with open(input_path, 'wb') as f:
                f.write(word_file.read())
            
            # Convert using LibreOffice
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(output_dir), str(input_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            # Find output file
            output_path = output_dir / f"{input_path.stem}.pdf"
            
            if not output_path.exists():
                raise Exception("Conversion failed")
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
            
            self._cleanup_files(output_path)
            
            filename = f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            return pdf_data, filename
            
        finally:
            self._cleanup_files(input_path)

    # ==================== PDF TO POWERPOINT ====================
    def pdf_to_pptx(self, pdf_file: io.BytesIO) -> Tuple[bytes, str]:
        """
        Convert PDF to PowerPoint.
        Converts each page to an image and adds to slides.
        """
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        pdf_doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        for page in pdf_doc:
            # Convert page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = pix.tobytes("png")
            
            # Add blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
            
            # Add image
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(
                img_stream, 
                Inches(0.25), Inches(0.25),
                Inches(9.5), Inches(7)
            )
        
        pdf_doc.close()
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        filename = f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        return output.read(), filename

    # ==================== POWERPOINT TO PDF ====================
    def pptx_to_pdf(self, pptx_file: io.BytesIO) -> Tuple[bytes, str]:
        """
        Convert PowerPoint to PDF.
        Uses LibreOffice for conversion.
        """
        input_path = self._get_temp_path('.pptx')
        output_dir = self.temp_dir
        
        try:
            with open(input_path, 'wb') as f:
                f.write(pptx_file.read())
            
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(output_dir), str(input_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            output_path = output_dir / f"{input_path.stem}.pdf"
            
            if not output_path.exists():
                raise Exception("Conversion failed")
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
            
            self._cleanup_files(output_path)
            
            filename = f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            return pdf_data, filename
            
        finally:
            self._cleanup_files(input_path)

    # ==================== EXCEL TO PDF ====================
    def excel_to_pdf(self, excel_file: io.BytesIO) -> Tuple[bytes, str]:
        """
        Convert Excel to PDF.
        Uses LibreOffice for conversion.
        """
        input_path = self._get_temp_path('.xlsx')
        output_dir = self.temp_dir
        
        try:
            with open(input_path, 'wb') as f:
                f.write(excel_file.read())
            
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', str(output_dir), str(input_path)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            output_path = output_dir / f"{input_path.stem}.pdf"
            
            if not output_path.exists():
                raise Exception("Conversion failed")
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
            
            self._cleanup_files(output_path)
            
            filename = f"converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            return pdf_data, filename
            
        finally:
            self._cleanup_files(input_path)


# Singleton instance
conversion_service = DocumentConversionService()
